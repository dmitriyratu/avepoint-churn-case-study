"""Build the AvePoint detailed-analysis deck as a .pptx."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

FIG = Path(r"C:\Users\dmitr\OneDrive\Dima\Job\Interview\AvePoint\outputs\figures")
OUT = Path(__file__).with_name("AvePoint_Detailed_Analysis.pptx")

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6B, 0x72, 0x80)
BLUE = RGBColor(0x2E, 0x5F, 0x8A)
RED = RGBColor(0xB0, 0x2E, 0x2E)
GREEN = RGBColor(0x1E, 0x7A, 0x4B)
RULE = RGBColor(0xD8, 0xDC, 0xE0)

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.75)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def _tb(slide, x, y, w, h):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _run(p, text, size, bold=False, color=INK, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return r


def header(slide, kicker, title, colour=INK):
    tf = _tb(slide, M, Inches(0.42), W - 2 * M, Inches(1.15))
    p = tf.paragraphs[0]
    _run(p, kicker.upper(), 12, bold=True, color=MUTED)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    _run(p2, title, 30, bold=True, color=colour)

    line = slide.shapes.add_shape(1, M, Inches(1.62), W - 2 * M, Emu(9525))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    line.shadow.inherit = False


def bullets(slide, items, top=Inches(1.95), width=None, size=17, left=M):
    tf = _tb(slide, left, top, width or (W - 2 * M), Inches(4.6))
    first = True
    for item in items:
        text, *style = item if isinstance(item, tuple) else (item,)
        colour = style[0] if style else INK
        bold = style[1] if len(style) > 1 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(11)
        if text.startswith("  "):
            p.level = 1
            _run(p, "– " + text.strip(), size - 2, color=MUTED)
        else:
            _run(p, text, size, bold=bold, color=colour)
    return tf


def picture(slide, name, top=Inches(1.9), height=Inches(4.5), left=None):
    pic = slide.shapes.add_picture(str(FIG / name), Inches(0), top, height=height)
    pic.left = left if left is not None else int((W - pic.width) / 2)
    return pic


def footnote(slide, text):
    tf = _tb(slide, M, H - Inches(0.72), W - 2 * M, Inches(0.4))
    _run(tf.paragraphs[0], text, 11, color=MUTED, italic=True)


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add(kicker, title, colour=INK):
    s = prs.slides.add_slide(BLANK)
    header(s, kicker, title, colour)
    return s


# ---------------------------------------------------------------- 1 title
s = prs.slides.add_slide(BLANK)
tf = _tb(s, M, Inches(2.3), W - 2 * M, Inches(2.6))
_run(tf.paragraphs[0], "SaaS Churn — Detailed Analysis", 44, bold=True)
p = tf.add_paragraph(); p.space_before = Pt(14)
_run(p, "RavenStack dataset · 500 accounts · why they leave, whether we can "
        "predict it, and what to do", 19, color=MUTED)
p = tf.add_paragraph(); p.space_before = Pt(26)
_run(p, "Headline: nothing about a customer predicts churn — but the date does.",
     20, bold=True, color=BLUE)
note(s, "The deck is organised around the three questions the product team asked. "
        "Two answers are negative; the negatives are the substance.")

# ---------------------------------------------------------------- 2 answers
s = add("The brief", "Three questions, three answers")
bullets(s, [
    ("1. Why are users leaving?", INK, True),
    ("  Not who — when. Churn hazard rose 2.8x per year (p = 2e-16) on a flat "
     "customer base. No attribute, segment, tenure band or stated reason "
     "explains who leaves."),
    ("2. Can we predict churn before it happens?", INK, True),
    ("  No. Three independent formulations land at chance: nested CV 0.534, "
     "Cox concordance 0.509, Cox global p = 0.57."),
    ("3. What actions improve retention?", INK, True),
    ("  Not identifiable — no intervention is recorded anywhere in the schema. "
     "And the economics say targeting is unnecessary."),
])
footnote(s, "Answers 1 and 2 are the same fact: a factor that moves every account "
            "at once leaves no cross-sectional variance for a model to fit.")

# ---------------------------------------------------------------- 3 framing
s = add("Framing", "The supplied label was not usable")
bullets(s, [
    ("accounts.churn_flag has no date attached and disagrees with the churn "
     "event log for 312 of 500 accounts — 37.6% agreement.", RED),
    "Replaced with the standard temporal formulation: given everything "
    "observable on 2024-06-30, will this account churn within 90 days?",
    "  Eligible: signed up before the cutoff, subscription still open, not "
    "already churned  ->  177 accounts",
    "  Label: first churn event inside the window  ->  54 positives (30.5%)",
    "  Features: only rows dated before the cutoff; fields resolving after it "
    "are censored",
    ("Found by counting, not by modelling. This is the single highest-leverage "
     "finding in the project.", BLUE, True),
])

# ---------------------------------------------------------------- 4 ladder
s = add("Question 2 — prediction", "Ten models, none clears chance")
picture(s, "04_model_ladder.png", height=Inches(4.3))
footnote(s, "Repeated stratified CV, 5 folds x 10 repeats, identical folds. "
            "No rung clears 0.5 at the lower bound of its interval.")

# ---------------------------------------------------------------- 5 nested
s = add("Question 2 — prediction", "Choosing the winner is itself a fitting step")
bullets(s, [
    ("Ladder maximum (optimistic):  0.583", MUTED),
    ("Nested CV (honest):  0.534 ± 0.016", INK, True),
    ("Chance:  0.500", MUTED),
    ("Cost of selecting the winner:  0.049 — most of the apparent signal.", RED, True),
    "Across 25 outer folds, eight of the ten rungs won at least once. A "
    "genuinely better model wins consistently; eight winners is what selecting "
    "on noise looks like.",
    "  Same effect measured independently: best-of-15 classifiers on shuffled "
    "labels scores 0.566, and the real winner's selection-corrected p is 0.300",
])

# ---------------------------------------------------------------- 6 lead time
s = add("Question 2 — prediction", "The signal does not survive a lead-time requirement")
picture(s, "03_horizon_buffer_sweep.png", height=Inches(4.0))
footnote(s, "Buffer = days of warning demanded. Every zero-buffer cell scores "
            "0.56-0.59; every cell demanding 30-60 days drops to 0.42-0.52. "
            "Accurate too late to act on.")

# ---------------------------------------------------------------- 7 THE finding
s = add("Question 1 — why", "The one large effect: calendar time", BLUE)
picture(s, "12_calendar_hazard.png", height=Inches(4.2))
footnote(s, "Poisson trend with log(at-risk) offset: rate ratio 1.089/month, "
            "95% CI [1.067, 1.112], x2.8 per year, p = 2e-16. At-risk base flat "
            "at ~200 accounts, so this is not a growth artefact.")
note(s, "This is the payoff slide. The same-sized customer base leaves four times "
        "faster at the end of the window than at the start.")

# ---------------------------------------------------------------- 8 trap 1
s = add("Question 1 — why", "A trap: the hazard falls sharply with tenure")
picture(s, "12_hazard_shape.png", height=Inches(4.0))
footnote(s, "Weibull rho = 0.737, 95% CI [0.673, 0.801], p = 1.7e-13 — the most "
            "emphatic result anywhere in the project. Read alone: churn is "
            "front-loaded, so fix onboarding.")

# ---------------------------------------------------------------- 9 trap 2
s = add("Question 1 — why", "...and it disappears under decomposition", RED)
picture(s, "12_cohort_gradient.png", height=Inches(3.9))
footnote(s, "Within each signup cohort rho returns to 1 (0.87, 1.25, 1.17, 1.06, "
            "0.88 — none significant). At fixed tenure, 30-day survival runs "
            "0.98 (2023Q2) to 0.59 (2024Q4). Recent cohorts churn faster at every "
            "age and supply most of the short tenures.")

# ---------------------------------------------------------------- 10 reversal
s = add("Consequence", "This reverses a recommendation", RED)
bullets(s, [
    ("Earlier in this project: \"onboard the first 6 months harder\" — resting "
     "on days_since_signup being the top single feature and on a tenure-band "
     "table.", MUTED),
    ("Both are the same composition artefact seen through a weaker lens.", RED, True),
    "Within any signup cohort, churn is memoryless: an account at day 300 is "
    "as likely to leave next week as one at day 10.",
    ("Do not build the structured onboarding programme.", RED, True),
    "The earlier recommendation is left in the repo rather than quietly "
    "deleted, with the reversal recorded where it was found. The sequence is "
    "the lesson.",
])

# ---------------------------------------------------------------- 11 reasons
s = add("Question 1 — why", "The field that records \"why\" is noise")
bullets(s, [
    "churn_events carries a reason_code and free-text comment for every "
    "departure. It fails three independent tests:",
    ("  Distribution over 6 codes vs uniform:  chi-square p = 0.70  (no dominant cause)"),
    ("  Reason vs 9 behavioural measures, BH-corrected:  min p = 0.91"),
    ("  Reason code vs the customer's own free text:  Cramer's V = 0.09"),
    ("Accounts coded \"support\" have the same ticket counts and satisfaction "
     "scores as accounts coded \"pricing\". Accounts coded \"features\" used the "
     "same number of features as everyone else.", RED),
    ("A \"top churn reasons\" chart from this field would be confident and "
     "entirely fictional — the most likely wrong output of this dataset.",
     BLUE, True),
])

# ---------------------------------------------------------------- 12 shap 1
s = add("Question 1 — why", "The driver chart that would go in the deck")
picture(s, "13_shap_importance.png", height=Inches(4.2))
footnote(s, "TreeSHAP on the gradient booster. Reads as: support load per seat "
            "drives churn, then contract size, then tenure. All three are "
            "plausible, which is the problem.")

# ---------------------------------------------------------------- 13 shap 2
s = add("Question 1 — why", "The same chart from shuffled labels", RED)
picture(s, "13_shap_null.png", height=Inches(3.3), top=Inches(1.95))
bullets(s, [
    ("Top feature p = 0.24  ·  12 different features win across 25 bootstrap "
     "resamples  ·  permutation importance on held-out data is beaten by noise "
     "(p = 1.0)  ·  TreeSHAP and the logistic model share 2 of their top 10.",
     RED, True),
], top=Inches(5.45), size=14)
footnote(s, "Explainability tooling has no failure mode: SHAP does not decline "
            "to rank when the model is worthless, and no library warns you.")

# ---------------------------------------------------------------- 14 causal
s = add("Question 3 — actions", "No intervention exists, so nothing is identified")
picture(s, "14_ate_forest.png", height=Inches(3.4), top=Inches(1.9))
bullets(s, [
    ("Cross-fitted AIPW on four observational proxies. Every interval crosses "
     "zero; upgrade is closest at ATE -0.110, p = 0.14. Adjustment already ate "
     "part of the naive -0.126 — upgraders differ from non-upgraders by SMD "
     "0.87 on subscription count, 0.81 on usage, 0.73 on tenure.", INK),
], top=Inches(5.45), size=14)

# ---------------------------------------------------------------- 15 evalue
s = add("Question 3 — actions", "Two stress tests that end the discussion")
bullets(s, [
    ("E-values 1.3 to 2.2.", INK, True),
    "  A confounder associated with both upgrading and retention by a risk "
    "ratio of 2.2 erases the whole estimate. \"The account was doing well\" "
    "clears that easily — and is unmeasurable in this schema.",
    ("Placebo band ±15pp.", INK, True),
    "  Running the entire AIPW pipeline on randomly assigned treatments, where "
    "the true effect is exactly zero, produces estimates up to ±0.15. Every "
    "observed estimate sits inside it.",
    ("An independent closed-form power calculation puts the minimum detectable "
     "effect at 17.2%. Two methods, two points apart.", BLUE, True),
])
footnote(s, "Uplift modelling does not rescue it: Qini p = 0.30, and treatment "
            "was self-selected, so ranked uplift mixes effect with selection.")

# ---------------------------------------------------------------- 16 economics
s = add("Question 3 — actions", "The arithmetic that should have come first", GREEN)
picture(s, "15_breakeven_grid.png", height=Inches(3.5), top=Inches(1.9))
bullets(s, [
    ("Break-even churn probability 10.3%  ·  cohort base rate 30.5%  ->  "
     "contact everyone. A ranking cannot improve on a policy that is correct "
     "for everyone.", GREEN, True),
], top=Inches(5.55), size=15)
footnote(s, "CLV $7,310 integrated from the measured survival curve. Model best "
            "threshold nets $53,000 against $52,400 for treat-all — a 1% gain, "
            "itself optimistic from in-sample threshold selection.")

# ---------------------------------------------------------------- 17 power
s = add("Question 3 — actions", "Only a churn-halving effect is testable here")
picture(s, "15_experiment_power.png", height=Inches(4.0))
footnote(s, "At ~21 signups/month: 15pp reads out in ~15 months; 5pp needs 2,527 "
            "accounts and over ten years. A fact about having 500 customers, not "
            "about the analysis.")

# ---------------------------------------------------------------- 18 actions
s = add("Recommendations", "In order of leverage")
bullets(s, [
    ("1.  Find what changed between 2023 and 2024.", INK, True),
    "  Hazard x2.8/yr on a flat base. An order of magnitude larger than "
    "anything account-level. The extract holds no calendar-varying covariate — "
    "no pricing log, no release history, no competitor events.",
    ("2.  Fix the churn label and the usage timestamps.", INK, True),
    "  Three definitions agree on 20% of accounts; 19,128 of 24,979 usage rows "
    "predate their own subscription.",
    ("3.  Contact the whole at-risk cohort. Do not rank it.", INK, True),
    ("4.  Instrument interventions — every CSM touch, discount and campaign.", INK, True),
    "  Without this, question 3 stays unanswerable at any sample size.",
    ("5.  Do not build the onboarding programme.", RED, True),
])

# ---------------------------------------------------------------- 19 rigour
s = add("How this was kept honest", "Every claimed pattern has a null to beat")
bullets(s, [
    "Automated leakage suite gates every reported number — temporal "
    "provenance, forbidden columns by name, single-feature AUC, identifier "
    "leakage. It caught a leak code review missed.",
    "  Post-outcome columns are worth +0.37 AUC and land at a plausible 0.79, "
    "not an obviously broken 0.99 — which is what makes them dangerous",
    "51 tests. Three pin bugs that were live and produced plausible output:",
    "  a propensity trim that silently trimmed nothing; an AIPW estimate that "
    "left [0,1] (a -0.17 churn rate); a KM interval that disagreed with its own "
    "point estimate",
    ("Nulls throughout: shuffled-label SHAP, selection-corrected p for the "
     "classifier sweep, placebo treatments for AIPW, permutation null for "
     "segment gaps, positive controls for the pipeline itself.", BLUE),
])
footnote(s, "15 notebooks, ~26 min end to end, all regenerated from source.")

prs.save(OUT)
print(f"{OUT}  ({OUT.stat().st_size/1e6:.2f} MB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
