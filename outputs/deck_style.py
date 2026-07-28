"""Shared layout and styling for both decks.

Language rule for everything built on top of this: each line is a complete
sentence in short plain words. The methods are already hard. The English
should not be.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

FIG = Path(r"C:\Users\dmitr\OneDrive\Dima\Job\Interview\AvePoint\outputs\figures")

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x62, 0x70)
BLUE = RGBColor(0x2E, 0x5F, 0x8A)
RED = RGBColor(0xB0, 0x2E, 0x2E)
GREEN = RGBColor(0x1E, 0x7A, 0x4B)
RULE = RGBColor(0xD8, 0xDC, 0xE0)

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.75)


def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def tb(slide, x, y, w, h):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    return tf


def run(p, text, size, bold=False, color=INK, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return r


def header(slide, kicker, title, colour=INK, title_size=28):
    tf = tb(slide, M, Inches(0.40), W - 2 * M, Inches(1.15))
    run(tf.paragraphs[0], kicker.upper(), 12, bold=True, color=MUTED)
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    run(p, title, title_size, bold=True, color=colour)
    line = slide.shapes.add_shape(1, M, Inches(1.56), W - 2 * M, Emu(9525))
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    line.shadow.inherit = False


def bullets(slide, items, top=Inches(1.88), width=None, size=16, left=M,
            height=Inches(5.0), space=9):
    """Items are text, or (text, colour), or (text, colour, bold).

    Two leading spaces marks a sub-point.
    """
    tf = tb(slide, left, top, width or (W - 2 * M), height)
    first = True
    for item in items:
        text, *style = item if isinstance(item, tuple) else (item,)
        colour = style[0] if style else INK
        bold = style[1] if len(style) > 1 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(space)
        if text.startswith("  "):
            p.level = 1
            run(p, text.strip(), size - 2, color=colour if style else MUTED)
        elif not text:
            run(p, " ", 6)
        else:
            run(p, text, size, bold=bold, color=colour)
    return tf


def picture(slide, name, top=Inches(1.85), height=Inches(4.3), left=None):
    pic = slide.shapes.add_picture(str(FIG / name), Inches(0), top, height=height)
    pic.left = left if left is not None else int((W - pic.width) / 2)
    return pic


def footnote(slide, text, colour=MUTED):
    tf = tb(slide, M, H - Inches(0.80), W - 2 * M, Inches(0.55))
    run(tf.paragraphs[0], text, 11.5, color=colour, italic=True)


def note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def stat(slide, x, y, value, label, colour=INK, w=Inches(3.0)):
    """A large number with a caption under it, for the executive deck."""
    tf = tb(slide, x, y, w, Inches(1.5))
    run(tf.paragraphs[0], value, 40, bold=True, color=colour)
    p = tf.add_paragraph()
    p.space_before = Pt(2)
    run(p, label, 13, color=MUTED)
    return tf
