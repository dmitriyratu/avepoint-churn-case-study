"""Append the customer-profile slide to the staging deck in Downloads.

Re-runnable: a slide already carrying this title is dropped first.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))

from add_slide_data_model import drop_existing  # noqa: E402
from deck_style import FIG, INK, MUTED, W, footnote, header, note  # noqa: E402

DECK = Path.home() / "Downloads" / "AvePoint_New_Slides.pptx"
FIGURE = FIG / "20_customer_profile.png"
TITLE = "What a customer looks like in this file"


def main():
    prs = Presentation(str(DECK))
    drop_existing(prs, TITLE)
    before = len(prs.slides)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Part 2 · Data understanding", TITLE, INK)

    pic = s.shapes.add_picture(str(FIGURE), Inches(0), Inches(1.55),
                               height=Inches(4.55))
    pic.left = int((W - pic.width) / 2)

    footnote(s, "The example is the median customer by contract count, not a "
                "worst case. Most of the file looks like this.", MUTED)

    note(s, "The slide that answers \"what does one row actually mean here\".\n\n"
            "Read the left panel first. Ten contracts, all open, tier moving "
            "Basic to Pro to Enterprise and back with no upgrade or downgrade "
            "flags set, seats disagreeing with the account record, and a total "
            "monthly bill of $33,506 on an account that says 28 seats and one "
            "plan.\n\n"
            "435 of 500 customers hold all three tiers at once. So there is no "
            "such thing as 'the customer's plan' or 'the customer's price' in "
            "this file — every spend feature needs a rule, and I use the most "
            "recent contract and the average across contracts.\n\n"
            "The green box is the part that matters for credibility. Price per "
            "seat is exactly $19, $49 and $199 by tier with no exceptions, and "
            "ARR is always twelve times MRR. Each row was built by a rule, "
            "carefully. That is why subscriptions is the table I lean on, and "
            "why the failure elsewhere is about relationships between rows "
            "rather than sloppiness inside them.")

    prs.save(str(DECK))
    print(f"{DECK.name}: {before} + 1 -> {len(prs.slides)} slides")
    print(f"  picture {pic.width / 914400:.2f} x {pic.height / 914400:.2f} in")


if __name__ == "__main__":
    main()
