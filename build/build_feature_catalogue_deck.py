"""A standalone one-slide deck with the feature catalogue.

Writes its own file, so nothing else can be disturbed.

Run from the repo root:  python build/build_feature_catalogue_deck.py
"""
import sys
from pathlib import Path

from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))

from deck_style import (FIG, INK, W, blank, footnote, header,  # noqa: E402
                        new_deck, note)

OUT = Path.home() / "Downloads" / "AvePoint_Feature_Catalogue.pptx"
FIGURE = FIG / "22_feature_catalogue.png"

prs = new_deck()
s = blank(prs)
header(s, "Part 2 · Feature engineering", "73 features, by transformation", INK)

# Sized by width: the table is wider than it is tall, so height is what falls
# out, not what is chosen.
pic = s.shapes.add_picture(str(FIGURE), Inches(0), Inches(1.95),
                           width=Inches(12.10))
pic.left = int((W - pic.width) / 2)

footnote(s, "Every feature appears exactly once. The build checks this against "
            "the pipeline, so the table cannot go stale.")

note(s, "For \"walk me through your feature engineering\".\n\n"
        "Grouping by source table says where a feature came from. Grouping by "
        "technique says what I did. Nine techniques, 73 features, each listed "
        "once.\n\n"
        "Definitions if pushed. Trailing windows are 30, 60, 90 and 180 days. "
        "Acceleration divides a short window rate by a long window rate, each "
        "over its own length, so it reads 1.0 when nothing has changed. The "
        "slope is least squares on weekly counts over 180 days — it catches a "
        "steady decline that window ratios miss. Rhythm is the mean and max "
        "gap between active days.\n\n"
        "Four choices worth defending:\n"
        "Missing means something different per family. A count of zero is a "
        "real zero. A recency of never is maximally stale, not today, so it "
        "takes the window length. An unknown rate stays blank and is imputed "
        "inside the fold.\n"
        "Per-seat divides by the newest contract that had already started, not "
        "by accounts.seats — that column is current as of extraction and would "
        "leak a later state.\n"
        "Six near-duplicates were dropped above 0.98 correlation, including "
        "feature_breadth against unique_features_used.\n"
        "A ticket still open at the cutoff has no resolution time and no "
        "satisfaction score. An automatic check found the five this hit.\n\n"
        "The same function builds training and production features. Training "
        "passes the cutoff, production passes today.\n\n"
        "What I would add with better data: session interarrival times, seats "
        "used against seats bought, and how many distinct people at the account "
        "touch the product. None is recoverable from this extract.")

prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides)} slide)")
print(f"  picture {pic.width / 914400:.2f} x {pic.height / 914400:.2f} in")
