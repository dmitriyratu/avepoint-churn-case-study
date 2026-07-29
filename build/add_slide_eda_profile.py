"""Append the univariate EDA slide to the staging deck in Downloads.

Re-runnable: a slide already carrying this title is dropped first.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))

from add_slide_data_model import drop_existing  # noqa: E402
from deck_style import (FIG, INK, RED, W, bullets, footnote,  # noqa: E402
                        header, note)

DECK = Path.home() / "Downloads" / "AvePoint_New_Slides.pptx"
FIGURE = FIG / "21_eda_profile.png"
TITLE = "Every field, one at a time"


def main():
    prs = Presentation(str(DECK))
    drop_existing(prs, TITLE)
    before = len(prs.slides)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Part 2 · Exploratory analysis", TITLE, INK)

    pic = s.shapes.add_picture(str(FIGURE), Inches(0), Inches(1.34),
                               height=Inches(4.30))
    pic.left = int((W - pic.width) / 2)

    bullets(s, [
        ("Spend and seats are right-skewed, the way real billing data is — a "
         "long tail of large accounts, so I log them before modelling. Tickets "
         "and usage per customer are not: a narrow band with no tail, which is "
         "what a fixed draw per customer looks like rather than behaviour.",
         INK, True),
        ("Every categorical sits on the even-split line. Ticket priority "
         "splits 26 / 26 / 25 / 24 across four levels, and churn reason six "
         "ways within four points. Real triage is never uniform — urgent is "
         "rare — and real churn has a dominant reason.", RED, True),
    ], top=Inches(5.78), size=12.5, space=4)

    footnote(s, "Missingness is recorded, never imputed. A ticket with no "
                "satisfaction score has not been rated, and that is itself a "
                "feature.")

    note(s, "The conventional univariate pass, in the conventional order: "
            "distributions, then categorical frequencies, then missingness.\n\n"
            "What to point at. Spend and seats look like real billing data — "
            "right-skewed, long tail, median $931 and 15 seats. Tickets and "
            "usage do not: tickets run 1 to 11 with a median of 4, usage 10 to "
            "101 with a median of 50. No power law, no heavy users, no dormant "
            "accounts. That shape is a per-customer draw, not behaviour.\n\n"
            "The dashed line on each categorical is what an even split would "
            "be. Plan, industry and referral source sit near it, which is "
            "unremarkable on 500 customers. Ticket priority and churn reason "
            "sit on it exactly, and those are the two that should not — urgent "
            "tickets are rare in every support organisation, and churn reasons "
            "always have a dominant category.\n\n"
            "Missingness: only three columns have any. end_date is missing for "
            "90% because contracts never end. satisfaction_score is missing for "
            "41% and takes only 3, 4 and 5 when present, so both the value and "
            "its absence are degenerate. Nothing is imputed — a missing "
            "satisfaction score means the ticket was not rated, which is "
            "information, and filling it would invent a rating.")

    prs.save(str(DECK))
    print(f"{DECK.name}: {before} + 1 -> {len(prs.slides)} slides")
    print(f"  picture {pic.width / 914400:.2f} x {pic.height / 914400:.2f} in")


if __name__ == "__main__":
    main()
