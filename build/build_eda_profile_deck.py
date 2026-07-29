"""A standalone one-slide deck with the univariate EDA panel.

Writes its own file rather than appending to anything, so nothing else can be
disturbed. Re-running overwrites only this deck.

Run from the repo root:  python build/build_eda_profile_deck.py
"""
import sys
from pathlib import Path

from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))

from deck_style import (FIG, INK, RED, W, blank, bullets, footnote,  # noqa: E402
                        header, new_deck, note)

OUT = Path.home() / "Downloads" / "AvePoint_EDA_Profile.pptx"
FIGURE = FIG / "21_eda_profile.png"

prs = new_deck()
s = blank(prs)
header(s, "Part 2 · Exploratory analysis", "Every field, one at a time", INK)

pic = s.shapes.add_picture(str(FIGURE), Inches(0), Inches(1.34),
                           height=Inches(4.30))
pic.left = int((W - pic.width) / 2)

bullets(s, [
    ("Spend and seats behave like real billing data. A few very large accounts "
     "stretch the scale, so I take logs before modelling.", INK, True),
    ("Tickets and usage do not behave like activity. Every customer has between "
     "1 and 11 tickets and between 10 and 101 usage rows. Real products have "
     "power users and dormant accounts. This one has neither, so I treat volume "
     "as a count and never as engagement.", RED, True),
    ("Every category is close to an even split. Ticket priority runs "
     "26 / 26 / 25 / 24, churn reasons 19 / 17 / 17 / 16 / 15 / 15. Urgent "
     "tickets are rare in any real support desk and real churn has one dominant "
     "reason, so none of these is a segment worth targeting.", RED, True),
], top=Inches(5.62), size=12.5, space=3)

footnote(s, "Blanks are kept as blanks. A ticket with no satisfaction score "
            "was never rated, so filling one in would invent a rating. Only "
            "three columns have any missing values at all.")

note(s, "The conventional univariate pass, in the conventional order: "
        "distributions, then categorical frequencies, then missingness. Every "
        "panel names the table and column it came from.\n\n"
        "What to point at. Spend and seats look like real billing data — "
        "right-skewed, long tail, median $931 and 15 seats. Tickets and usage "
        "do not: tickets run 1 to 11 with a median of 4, usage 10 to 101 with a "
        "median of 50. No power law, no heavy users, no dormant accounts. That "
        "shape is a per-customer draw, not behaviour.\n\n"
        "The dashed line on each categorical is what an even split would be. "
        "Plan, industry and referral source sit near it, which is unremarkable "
        "on 500 customers. Ticket priority and churn reason sit on it exactly, "
        "and those are the two that should not — urgent tickets are rare in "
        "every support organisation, and churn reasons always have a dominant "
        "category.\n\n"
        "Missingness: only three columns have any. end_date is missing for 90% "
        "because contracts never end. satisfaction_score is missing for 41% and "
        "takes only 3, 4 and 5 when present, so both the value and its absence "
        "are degenerate. Nothing is imputed — a missing satisfaction score "
        "means the ticket was not rated, which is information, and filling it "
        "would invent a rating.")

prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides)} slide)")
print(f"  picture {pic.width / 914400:.2f} x {pic.height / 914400:.2f} in")
