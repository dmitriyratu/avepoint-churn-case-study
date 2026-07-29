"""Append the data-model slide to the staging deck in Downloads.

Re-runnable: a slide already carrying this title is dropped first, so editing
the figure and running again replaces it rather than stacking duplicates.
Nothing else in the deck is touched.
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

sys.path.insert(0, str(Path(__file__).resolve().parent))

from deck_style import FIG, INK, W, header, note  # noqa: E402

DECK = Path.home() / "Downloads" / "AvePoint_New_Slides.pptx"
FIGURE = FIG / "19_data_model.png"
TITLE = "Five tables, and which of them can be trusted"


def drop_existing(prs, title):
    """Remove any slide already carrying this title, so a rerun replaces it."""
    listing = prs.slides._sldIdLst
    for slide in list(prs.slides):
        text = " ".join(sh.text_frame.text for sh in slide.shapes
                        if sh.has_text_frame)
        if title not in text:
            continue
        for entry in list(listing):
            if prs.part.related_part(entry.rId) is slide.part:
                listing.remove(entry)
                prs.part.drop_rel(entry.rId)


def main():
    prs = Presentation(str(DECK))
    drop_existing(prs, TITLE)
    before = len(prs.slides)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, "Part 2 · Data exploration", TITLE, INK)

    # The figure carries the read-out as well as the diagram, so it gets the
    # whole slide below the header and there is nothing to say underneath it.
    pic = s.shapes.add_picture(str(FIGURE), Inches(0), Inches(1.40),
                               height=Inches(5.90))
    pic.left = int((W - pic.width) / 2)

    note(s, "The slide to open with if anyone asks what the data looks like.\n\n"
            "Read the arrows, not the boxes. Referential integrity is perfect: "
            "every subscription_id in feature_usage resolves, every account_id "
            "in support_tickets and churn_events resolves, and no contract "
            "predates its account. A schema check passes cleanly, which is why "
            "a cleaning pipeline does not catch this — and the other published "
            "analysis of this dataset ran one and reported nothing.\n\n"
            "What fails is temporal integrity, and nothing in a data dictionary "
            "tests it. 13,198 of 25,000 usage rows and 1,077 of 2,000 tickets "
            "are dated before their customer signed up. Correlation between "
            "event date and signup date is 0.002 and 0.016. Every behavioural "
            "feature in the model is built from those two tables.\n\n"
            "The left column of the read-out is every check a reviewer would "
            "ask for. The right column is why passing them settles nothing.")

    prs.save(str(DECK))
    print(f"{DECK.name}: {before} + 1 -> {len(prs.slides)} slides")
    print(f"  picture {pic.width / 914400:.2f} x {pic.height / 914400:.2f} in")


if __name__ == "__main__":
    main()
